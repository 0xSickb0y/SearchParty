# -*- coding: utf-8 -*-

import re
import csv
import time
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
from email import message_from_file
from openpyxl import load_workbook
from pdfminer.high_level import extract_text
from .Utils.ColorFormatting import display_info
from .Utils.CustomExceptions import NoDataFound

wordlists = {}
error_log = {}
data_found = {}
data_indexes = {}

regex_patterns = {"CPF": re.compile(r'\b\d{9}/\d{2}\b|\b\d{3}\.\d{3}\.\d{3}-\d{2}\b'),
                  "RG": re.compile(r'\b\d{2}\.\d{3}\.\d{3}-\d{1}\b'),
                  "Email Addresses": re.compile(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'),
                  "Phone Numbers": re.compile(r'(?:\+\d{1,3}\s)?(?:\(\d{2}\)|\d{2})\s?\b\d{5}-\d{4}\b'),
                  "Employee IDs": re.compile(r'\d{3}\.\d{5}\.\d{2}-\d'),
                  "National Health Card": re.compile(r'\d{3}[ .-]\d{4}[ .-]\d{4}[ .-]\d{4}')}

wordlist_paths = ['wordlists/Ethnic Groups',
                  'wordlists/Financial Information',
                  'wordlists/Gender Orientation',
                  'wordlists/Legal Information',
                  'wordlists/Medical Records',
                  'wordlists/Political Preferences',
                  'wordlists/Property Information',
                  'wordlists/Religion and Faith',
                  'wordlists/Travel History',
                  'wordlists/Vehicle Documents']


def process_wordlists():
    for file_path in wordlist_paths:
        with open(file_path, 'r') as file:
            wordlist = file_path.split('wordlists/')[1]
            lines = [line.strip() for line in file.readlines()]
            wordlists[wordlist] = lines


def extract_from_pdf(args, pdf_path, tesseract_path=None):
    print("Extracting text from PDF files")
    for pdf_file in pdf_path:
        try:
            text = extract_text(pdf_file)
            find_patterns(args, text, pdf_file)
        except Exception:
            if 'pdf' not in error_log:
                error_log['pdf'] = {'count': 1}
            else:
                error_log['pdf']['count'] += 1


def extract_from_pptx(args, pptx_path, tesseract_path=None):
    print("Extracting text from PowerPoint files")
    for pptx_file in pptx_path:
        try:
            text = ''
            presentation = Presentation(pptx_file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        for line in shape.text.split('\n'):
                            text += f'{line}\n'
            find_patterns(args, text, pptx_file)
        except Exception:
            if 'pptx' not in error_log:
                error_log['pptx'] = {'count': 1}
            else:
                error_log['pptx']['count'] += 1


def extract_from_docx(args, docx_path, tesseract_path=None):
    print("Extracting text from Word Documents")
    for docx_file in docx_path:
        try:
            text = ''
            wdoc = Document(docx_file)
            for wdoc_p in wdoc.paragraphs:
                for line in wdoc_p.text.split('\n'):
                    text += f'{line}\n'
            find_patterns(args, text, docx_file)
        except Exception:
            if 'docx' not in error_log:
                error_log['docx'] = {'count': 1}
            else:
                error_log['docx']['count'] += 1


def extract_from_xlsx(args, xlsx_path, tesseract_path=None):
    print("Extracting text from Spreadsheets")
    for xlsx_file in xlsx_path:
        try:
            text = ''
            wbook = load_workbook(filename=xlsx_file)
            wsheet = wbook.active
            for row in wsheet.iter_rows(values_only=True):
                for cell in row:
                    for line in str(cell).split('\n'):
                        text += f'{line}\n'
            find_patterns(args, text, xlsx_file)
        except Exception:
            if 'xlsx' not in error_log:
                error_log['xlsx'] = {'count': 1}
            else:
                error_log['xlsx']['count'] += 1


def extract_from_txt(args, txt_path, tesseract_path=None):
    print("Extracting text from Text files")
    for txt_file in txt_path:
        try:
            with open(txt_file, 'r', encoding='utf-8') as file:
                text = file.read()
                find_patterns(args, text.strip(), txt_file)
        except Exception:
            if 'txt' not in error_log:
                error_log['txt'] = {'count': 1}
            else:
                error_log['txt']['count'] += 1


def extract_from_csv(args, csv_path, tesseract_path=None):
    print("Extracting text from CSV files")
    for csv_file in csv_path:
        try:
            text = ''
            with open(csv_file, 'r', encoding='utf-8') as file:
                reader = csv.reader(file)
                for row in reader:
                    text += str(row)
            find_patterns(args, text, csv_file)
        except Exception:
            if 'csv' not in error_log:
                error_log['csv'] = {'count': 1}
            else:
                error_log['csv']['count'] += 1


def extract_from_jpeg(args, jpeg_path, tesseract_path=None):
    pytesseract.tesseract_cmd = tesseract_path
    print("Extracting text from JPEG images")
    for jpeg_file in jpeg_path:
        try:
            text = ''
            for line in pytesseract.image_to_string(jpeg_file).split('\n'):
                text += f'{line}\n'
            find_patterns(args, text, jpeg_file)
        except Exception:
            if 'jpeg' not in error_log:
                error_log['jpeg'] = {'count': 1}
            else:
                error_log['jpeg']['count'] += 1


def extract_from_png(args, png_path, tesseract_path=None):
    pytesseract.tesseract_cmd = tesseract_path
    print("Extracting text from PNG images")
    for png_file in png_path:
        try:
            text = ''
            for line in pytesseract.image_to_string(png_file).split('\n'):
                text += f'{line}\n'
            find_patterns(args, text, png_file)
        except Exception:
            if 'png' not in error_log:
                error_log['png'] = {'count': 1}
            else:
                error_log['png']['count'] += 1


def extract_from_bmp(args, bmp_path, tesseract_path=None):
    pytesseract.tesseract_cmd = tesseract_path
    print("Extracting text from Bitmap images")
    for bmp_file in bmp_path:
        try:
            text = ''
            for line in pytesseract.image_to_string(bmp_file).split('\n'):
                text += f'{line}\n'
            find_patterns(args, text, bmp_file)
        except Exception:
            if 'bmp' not in error_log:
                error_log['bmp'] = {'count': 1}
            else:
                error_log['bmp']['count'] += 1


def extract_from_tiff(args, tiff_path, tesseract_path=None):
    pytesseract.tesseract_cmd = tesseract_path
    print("Extracting text from TIFF images")
    for tiff_file in tiff_path:
        try:
            text = ''
            for line in pytesseract.image_to_string(tiff_file).split('\n'):
                text += f'{line}\n'
            find_patterns(args, text, tiff_file)
        except Exception:
            if 'tiff' not in error_log:
                error_log['tiff'] = {'count': 1}
            else:
                error_log['tiff']['count'] += 1


def extract_from_webp(args, webp_path, tesseract_path=None):
    pytesseract.tesseract_cmd = tesseract_path
    print("Extracting text from WebP images")
    for webp_file in webp_path:
        try:
            text = ''
            for line in pytesseract.image_to_string(webp_file).split('\n'):
                text += f'{line}\n'
            find_patterns(args, text, webp_file)
        except Exception:
            if 'webp' not in error_log:
                error_log['webp'] = {'count': 1}
            else:
                error_log['webp']['count'] += 1


def extract_from_gif(args, gif_path, tesseract_path=None):
    pytesseract.tesseract_cmd = tesseract_path
    print("Extracting text from GIF images")
    for gif_file in gif_path:
        try:
            with Image.open(gif_file) as file:
                extracted_set = set()
                for frame in range(file.n_frames):
                    file.seek(frame)
                    frame_image = file.copy()
                    frame_text = pytesseract.image_to_string(frame_image)
                    for text in frame_text.split('\n'):
                        if text not in extracted_set:
                            extracted_set.add(text)
                find_patterns(args, text, gif_file)
        except Exception:
            if 'gif' not in error_log:
                error_log['gif'] = {'count': 1}
            else:
                error_log['gif']['count'] += 1


def extract_from_mail(args, email_path, tesseract_path=None):
    print("Extracting text from Email Messages")
    for email_file in email_path:
        try:
            with open(email_file, 'r', encoding='utf-8') as file:
                msg = message_from_file(file)
                text = ''
                for header in [msg[header] for header in msg.keys() if header.lower() in ['from', 'to', 'subject']]:
                    text += f'{header}\n'
                if msg.is_multipart():
                    for part in msg.walk():
                        if part.get_content_type() == 'text/plain':
                            payload = part.get_payload(decode=True)
                            for line in payload.splitlines():
                                line = line.decode(part.get_content_charset()).split("\n")
                                text += f'{line}\n'
                    find_patterns(args, text, email_file)
                else:
                    payload = msg.get_payload(decode=True)
                    for line in payload.splitlines():
                        text += f'{line.decode()}\n'
                    find_patterns(args, text, email_file)
        except Exception:
            if 'email' not in error_log:
                error_log['email'] = {'count': 1}
            else:
                error_log['email']['count'] += 1


def find_keywords(text):
    keywords_found = {}
    for category, keywords in wordlists.items():
        matches = []
        for keyword in keywords:
            if re.search(keyword, text, re.IGNORECASE):
                matches.append(keyword)
        if matches:
            keywords_found[category] = matches
    return keywords_found


def find_patterns(args, text, path):

    keywords_found = find_keywords(text)

    if args.findme:
        for argument in args.findme:
            key = ''.join(argument.split())
            if any(value.search(text) for value in [regex_patterns[key]]):
                if argument not in data_found:
                    data_found[argument] = []
                info = {'file': path}
                if any(value for value in keywords_found.values()):
                    info['keywords'] = keywords_found
                data_found[argument].append(info)
                data_indexes[argument] = data_indexes.get(argument, 0) + 1
    else:
        for pattern_name, pattern in regex_patterns.items():
            for match in pattern.finditer(text):
                data_indexes[pattern_name] = data_indexes.get(pattern_name, 0) + 1
                if pattern_name not in data_found:
                    data_found[pattern_name] = set()
                data_found[pattern_name].add(path)
                for key, value in ((key, value) for key, value in keywords_found.items() if value):
                    if key not in data_found:
                        data_found[key] = set()
                    if key not in data_indexes:
                        data_indexes[key] = len(value)
                    data_found[key].add(path)


def process_findme(args):
    if args.findme:
        args.findme = set(args.findme)
        regex_patterns.clear()
        for value in args.findme:
            value_key = ''.join(value.split())
            regex_patterns[f'{value_key}'] = re.compile(re.escape(value), re.IGNORECASE)


def process_extraction_methods(args, colors, file_types, data_filters, supported_files, tesseract_path=None):
    global regex_patterns

    print(display_info(f"Starting Data Extraction [{time.strftime('%H:%M:%S %d/%b/%Y')}]", colors))

    extraction_methods = {
        'bmp': extract_from_bmp, 'csv': extract_from_csv, 'gif': extract_from_gif,
        'pdf': extract_from_pdf, 'png': extract_from_png, 'txt': extract_from_txt,
        'docx': extract_from_docx, 'jpeg': extract_from_jpeg, 'pptx': extract_from_pptx,
        'tiff': extract_from_tiff, 'webp': extract_from_webp, 'xlsx': extract_from_xlsx,
        'mail': extract_from_mail
    }

    filtered_patterns = {}
    zip_keys = ['cpf', 'rg', 'email', 'phone', 'nit', 'cns']
    for zip_key, regex_key in zip(zip_keys, regex_patterns.keys()):
        if zip_key in data_filters:
            filtered_patterns[regex_key] = regex_patterns[regex_key]
    regex_patterns = filtered_patterns

    for file_type, method in extraction_methods.items():
        if file_type in file_types:
            if supported_files.get(file_type) is not None:
                method(args, supported_files[file_type], tesseract_path)

    if not data_found.keys():
        raise NoDataFound(args, colors)


def information_extractor(args, colors, file_types, data_filters, supported_files, tesseract_path=None):

    process_wordlists()
    process_findme(args)
    process_extraction_methods(args, colors, file_types, data_filters, supported_files, tesseract_path)
    return data_found, data_indexes, error_log
