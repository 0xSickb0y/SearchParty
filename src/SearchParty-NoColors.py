# -*- coding: utf-8 -*-

'''
https://github.com/0xSickb0y/SearchParty/
Made By: 0xSickb0y yuriVolp MarceloCostaM
v1.0
'''


import re
import os
import csv
import sys
import time
import psutil
import shutil
import sqlite3
import argparse
import platform
import datetime
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
from email import message_from_file
from magic import from_file
from filetype import guess_mime
from openpyxl import load_workbook
from pdfminer.high_level import extract_text


banner_str = r'''
   _____                      __       ____             __
  / ___/___  ____  ____ _____/ /_     / __ \____  _____/ /___  __
  \__ \/ _ \/ __ `/ ___/ ___/ __ \   / /_/ / __ `/ ___/ __/ / / /
 ___/ /  __/ /_/ / /  / /__/ / / /  / ____/ /_/ / /  / /_/ /_/ /
/____/\___/\__,_/_/   \___/_/ /_/  /_/    \__,_/_/   \__/\__  /
v1.0                                                     /___/

Ferramenta offline para mapeamento e análise de dados pessoais/sensíveis.
https://github.com/0xSickb0y/SearchParty/
'''

parser = argparse.ArgumentParser(
    prog='SearchParty.py',
    description=print(banner_str),
    epilog='')


def separate_args(arguments):
    return arguments.split(',')


parser.add_argument('-F', metavar="$file", dest='file', action='append', help='escanear arquivo')
parser.add_argument('-D', metavar="$directory", dest='directory', action='append', help='escanear diretório')
parser.add_argument('--find', metavar='', dest='findme', type=str, nargs='+', help="procurar valores específicos ('John Doe' '47.283.723-0')")
parser.add_argument('--loot', metavar="$name", dest='loot', type=separate_args, nargs='?', const=os.getcwd(), help='salvar resultados em um diretório (default: $current/loot)')
parser.add_argument('--database', metavar='$sql.db', dest='database', type=separate_args, nargs='?', const=os.getcwd(), help='salvar resultados em um banco de dados (default: $hostname.db)')
parser.add_argument('--data-type', metavar='$type', dest='datatype', type=separate_args, help='tipos de dados separados por vírgula')
parser.add_argument('--file-type', metavar='$type', dest='filetype', type=separate_args, help='tipos de arquivos separados por vírgula')
parser.add_argument('--copy-files', metavar='$dst', dest='copy', help='copiar arquivos para outro local')
parser.add_argument('--move-files', metavar='$dst', dest='move', help='mover arquivos para outro local')
parser.add_argument('--delete-files', dest='delete', action="store_true", help='excluir arquivos do sistema de arquivos')

args = parser.parse_args()


class SearchParty:

    def __init__(self):

        self.args = args

        self.regex_patterns = {"Cadastro de pessoa física": re.compile(r'\b\d{9}/\d{2}\b|\b\d{3}\.\d{3}\.\d{3}-\d{2}\b'),
                               "Registro geral": re.compile(r'\b\d{2}\.\d{3}\.\d{3}-\d{1}\b'),
                               "Endereço de email": re.compile(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'),
                               "Numero de telefone": re.compile(r'(?:\+\d{1,3}\s)?(?:\(\d{2}\)|\d{2})\s?\b\d{5}-\d{4}\b'),
                               "Numero de inscrição do trabalhador": re.compile(r'^\d{3}\.\d{5}\.\d{2}-\d$'),
                               "Cartão nacional de saúde": re.compile(r'^\d{3}[ .-]\d{4}[ .-]\d{4}[ .-]\d{4}$')}

        self.file_patterns = {'txt': 'text/plain',
                              'csv': 'text/csv',
                              'bmp': 'image/bmp',
                              'png': 'image/png',
                              'gif': 'image/gif',
                              'tiff': 'image/tiff',
                              'jpeg': 'image/jpeg',
                              'webp': 'image/webp',
                              'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                              'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                              'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                              'pdf': 'application/pdf',
                              'email': 'message/rfc822'}

        self.ethnic_groups = ['caucasiano', 'caucasiana', 'caucasianos', 'caucasianas', 'negro', 'negra', 'negros', 'negras', 'afrodescendente', 'afrodescendentes', 'hispanico', 'hispanica', 'hispanicos', 'hispanicas', 'latino', 'latina', 'latinos', 'latinas', 'asiático', 'asiática', 'asiáticos', 'asiáticas', 'nativoamericano', 'nativoamericana', 'nativoamericanos', 'nativoamericanas', 'indígena', 'indígenas', 'Maadans', "Ma'dan", 'árabe', 'árabes', 'pardo', 'parda', 'pardos', 'pardas', 'afro-latino', 'afro-latina', 'afro-latinos', 'afro-latinas', 'eslavo', 'eslava', 'eslavos', 'eslavas', 'polinésio', 'polinésia', 'polinésios', 'polinésias', 'melanésio', 'melanésia', 'melanésios', 'melanésias', 'micronésio', 'micronésia', 'micronésios', 'micronésias', 'berbere', 'cigano', 'cigana', 'ciganos', 'ciganas', 'romani', 'romanis', 'sinti', 'sintis', 'copta', 'coptas', 'assírio', 'assíria', 'assírios', 'assírias', 'druzo', 'druza', 'druzos', 'druzas', 'yazidi', 'yazidis', 'curdo', 'curda', 'curdos', 'curdas', 'brahmin', 'brahmina', 'brahmins', 'brahminas', 'dalit', 'dalita', 'dalits', 'dalitas', 'uigur', 'uigures', 'fulani', 'fulanis', 'igbo', 'igbos', 'yoruba', 'yorubas', 'sefardita', 'sefarditas', 'basco', 'bascos']
        self.financial_information = ['Programa de fidelidade', 'Número do cartão de fidelidade', 'Data de validade do cartão de crédito', 'Limite de crédito', 'Saldo disponível', 'Limite de saque', 'Histórico de compras', 'Produtos adquiridos', 'Valor das compras', 'Forma de pagamento', 'Data de vencimento', 'Histórico de transações financeiras', 'Descrição da transação', 'Valor da transação', 'Data da transação', 'Histórico de empréstimos', 'Valor do empréstimo', 'Taxa de juros', 'Valor das parcelas', 'Data de quitação', 'Histórico de investimentos', 'Tipo de investimento', 'Valor do investimento', 'Data de aplicação', 'Data de resgate', 'Rendimento', 'Número da Conta Bancária', 'Agência Bancária', 'Senha do Cartão', 'Código de Segurança', 'Extrato Bancário', 'Fatura do Cartão de Crédito', 'Comprovante de Pagamento', 'Comprovante de Renda', 'Declaração de Imposto de Renda', 'Número de Identificação Fiscal', 'Cheque Especial', 'Taxa de Câmbio', 'Saldo devedor', 'Tarifa Bancária', 'Investimentos', 'Rendimento', 'Valor de Resgate', 'Poupança', 'Certificado de Depósito Bancário', 'Valor de Depósito', 'Valor de Retirada', 'Seguro', 'Previdência Privada', 'Renda Fixa', 'Fundo de Investimento', 'Certificado de Depósito Bancário', 'Tesouro Direto', 'Bolsa de Valores', 'Dividendos', 'Renda Variável', 'Derivativos', 'Mercado Futuro', 'Índice Bovespa', 'Oferta Pública Inicial', 'Mercado de Capitais', 'Certificado de Depósito Interbancário', 'Letra de Crédito Imobiliário', 'Letra de Crédito do Agronegócio', 'Fundo Garantidor de Créditos', 'Índice Nacional de Preços ao Consumidor Amplo', 'Selic', 'Taxa Básica de Juros', 'Spread Bancário', 'Fluxo de Caixa', 'Conta Corrente', 'Crédito Pessoal', 'Hipoteca', 'Financiamento Imobiliário', 'Leasing', 'Amortização', 'Título de Capitalização', 'Aporte Financeiro', 'Conta Poupança', 'Adiantamento de Recebíveis', 'Dívida Pública', 'Nota Promissória', 'Imposto sobre Operações Financeiras', 'Juros Compostos', 'Juros Simples', 'Pagamento Mínimo do Cartão de Crédito', 'Spread de Crédito', 'Hedge', 'Debêntures', 'Garantia Real', 'Garantia Fiduciária', 'Valor Nominal', 'Valor Presente', 'Valor Futuro', 'Valor de Mercado', 'Patrimônio Líquido', 'Balanço Patrimonial', 'Demonstração de Resultado do Exercício', 'Lucro Líquido']
        self.gender_orientation = ['agenero', 'androgine', 'andrógino', 'aromantico', 'assexual', 'assexualidade', 'autossexual', 'autossexualidade', 'bigenero', 'bissexual', 'bissexualidade', 'cinzas-sexual', 'cinzassexual', 'cisgenero', 'cisgênero', 'demigenero', 'demigênero', 'demissexual', 'demissexualidade', 'fluidez de gênero', 'fluidez de genero', 'fluidez-de-genero', 'gay', 'genero fluido', 'gênero fluido', 'gênero-fluido', 'hetero', 'heteroafetivo', 'heterossexual', 'heterossexualidade', 'homoafetivo', 'homossexual', 'homossexualidade', 'hétero', 'intergenero', 'intergênero', 'lesbianismo', 'lesbica', 'litio romantico', 'litioromantico', 'lésbica', 'lítio romântico', 'lítioromântico', 'multigenero', 'multigênero', 'nao binario', 'nao-binario', 'naobinario', 'neutrois', 'não-binariedade', 'nãobinariedade', 'não-binário', 'não-binárie', 'omnissexual', 'omnissexualidade', 'pangenero', 'pangênero', 'pansexual', 'pansexualidade', 'poliamoroso', 'polissexual', 'polissexualidade', 'reciprosexual', 'reciprosexualidade', 'sapiosexual', 'sapiosexualidade', 'skoliosexual', 'skoliosexualidade', 'transgenero', 'transgênero', 'Acedisexual', 'Acedissexualidade', 'Agênero-fluxo', 'Aporagender', 'Aporagênero', 'Arofluxo', 'Aromantique', 'Arospécie', 'Autoandrogineidade', 'Autogynephilia', 'Autoqueer', 'Auto-sexo', 'Bigender-fluxo', 'Binaridade de gênero', 'Birromântico', 'Ceterosexual', 'Cinza-romântico', 'Cinzissexual', 'Cishomossexualidade', 'Demiromantique', 'Demivoid', 'Demisexualidade romântica', 'Espetrogênero', 'Espectralidade de gênero', 'Gênero-bissexual', 'Gênero-demissexual', 'Gênero-específico', 'Gênero fluido-fluxo', 'Gênero não-conforme', 'Gênero neutro-fluxo', 'Gênero-queer', 'Gênero questionador', 'Gênero-romântico', 'Gênero-fluido-romântico', 'Gênero-sui', 'Gênero-transespecífico', 'Gênero-transgressor', 'Gray-ace', 'Gray-aromântico', 'Heterorromântico', 'Homorromântico', 'Ignessexual', 'Intersexualidade de gênero', 'Lesbianismo-bissexual', 'Litio-queer', 'Multigênero-fluxo', 'Multissexualidade', 'Neutrois-romântico', 'Não-binaridade-fluxo', 'Nongênero-fluxo', 'Omnissexualidade romântica', 'Panfluxo', 'Panromantismo', 'Pansexuromântico', 'Plurissexualidade', 'Poliafeto', 'Polisexualidade-romântica', 'Queerflux', 'Recíproco-romântico', 'Reciprosexualidade romântica', 'Sedésico', 'Skaiofluxo', 'Skoliorromântico', 'Sui-romântico', 'Toric', 'Transmasculine', 'Transvestite', 'Ultrabinary', 'Xenogênero', 'Xenofluxo', 'Zedsexual', 'Zedsexuromântico']
        self.legal_information = ['Histórico de penalidades', 'Tipo de penalidade', 'Data da penalidade', 'Descrição da penalidade', 'Histórico de processos de licitação', 'Histórico de reclamações de consumidores', 'Empresa reclamada', 'Data da reclamação', 'Descrição da reclamação', 'Disputas judiciais', 'Parte demandante', 'Parte demandada', 'Data de abertura da disputa', 'Resultado da disputa', 'Histórico de registros em órgãos de proteção ao crédito', 'Órgão de proteção ao crédito', 'Data do registro', 'Descrição do registro', 'Audiência Preliminar', 'Audiência de Instrução e Julgamento', 'Contestação', 'Reconvenção', 'Embargos à Execução', 'Embargos de Terceiro', 'Medida Cautelar', 'Tutela Antecipada', 'Liminar', 'Agravo', 'Recurso Especial', 'Recurso Extraordinário', 'Recurso de Apelação', 'Recurso Adesivo', 'Recurso Inominado', 'Sentença', 'Acórdão', 'Mandado de Segurança', 'Habeas Corpus', 'Habeas Data', 'Ação Civil Pública', 'Ação Popular', 'Ação Rescisória', 'Mandado de Injunção', 'Ação Declaratória', 'Ação Cautelar', 'Ação Monitória', 'Ação Revisional', 'Ação de Busca e Apreensão', 'Ação de Despejo', 'Ação de Usucapião', 'Ação de Indenização', 'Ação de Execução', 'Ação de Divórcio', 'Ação de Alimentos', 'Ação de Investigação de Paternidade', 'Ação de Guarda', 'Ação de Interdição', 'Ação de Inventário', 'Ação de Arrolamento', 'Ação de Nunciação de Obra Nova', 'Ação de Desapropriação', 'Ação Popular', 'Ação Rescisória', 'Ação Direta de Inconstitucionalidade', 'Ação Direta de Inconstitucionalidade por Omissão', 'Arguição de Descumprimento de Preceito Fundamental', 'Mandado de Segurança Coletivo', 'Habeas Data Coletivo', 'Ação Civil Pública Ambiental', 'Ação de Improbidade Administrativa', 'Ação de Reparação de Danos', 'Ação de Usucapião Especial Urbana', 'Ação de Usucapião Especial Rural', 'Ação de Exoneração de Alimentos', 'Ação de Reconhecimento e Dissolução de União Estável', 'Ação de Regulamentação de Visitas', 'Ação de Responsabilidade Civil', 'Ação de Desconsideração da Personalidade Jurídica', 'Ação de Cobrança', 'Ação de Consignação em Pagamento', 'Ação de Anulação de Casamento', 'Ação de Separação Judicial', 'Ação de Anulação de Testamento', 'Ação de Depósito', 'Ação de Prestação de Contas', 'Ação de Rescisão Contratual', 'Ação de Nulidade de Cláusula Contratual', 'Ação de Despejo']
        self.medical_information = ['consultas médicas', 'Nome do médico', 'Especialidade médica', 'Diagnóstico', 'Prescrição médica', 'internações', 'Motivo da internação', 'Data de internação', 'Data de alta', 'Nome do hospital', 'Número do quarto', 'exames médicos', 'Tipo de exame', 'Resultado do exame', 'Data do exame', 'medicamentos', 'Nome do medicamento', 'Dosagem', 'Frequência', 'Data de início do tratamento', 'Data de término do tratamento', 'alergias', 'Tipo de alergia', 'Gravidade', 'vacinação', 'Tipo de vacina', 'Data da vacinação', 'Nome do vacinador', 'Local da vacinação', 'Informações médicas', 'Doenças crônicas', 'Saúde mental', 'Antecedentes familiares de doenças', 'cirurgias', 'Resultados de testes laboratoriais', 'transfusões sanguíneas', 'Regime alimentar', 'atividade física', 'consumo de álcool', 'consumo de tabaco', 'consumo de drogas ilícitas', 'Contato de emergência', 'Plano de saúde', 'Dados de seguro médico', 'Contatos de profissionais de saúde', 'peso corporal', 'pressão arterial', 'temperatura corporal', 'nível de glicose no sangue', 'nível de colesterol', 'batimentos cardíacos', 'sintomas', 'saúde familiar', 'saúde pessoal', 'cirurgias prévias', 'transfusões prévias', 'alergias alimentares', 'intolerâncias alimentares', 'dieta especial', 'histórico odontológico', 'próteses ou implantes', 'doenças infecciosas passadas', 'acidentes', 'fraturas', 'tratamentos psicológicos ou psiquiátricos anteriores', 'condições de saúde ocupacional', 'exposição a substâncias tóxicas ou químicas', 'exames de imagem', 'ecografias', 'tomografias', 'ressonâncias magnéticas', 'radiografias', 'mamografias', 'densitometrias ósseas', 'ultrassonografias', 'eletrocardiogramas (ECG)', 'eletroencefalogramas (EEG)', 'monitoramento cardíaco Holter', 'monitoramento da pressão arterial ambulatorial (MAPA)', 'espirometrias', 'testes de função pulmonar', 'polissonografias', 'biópsias', 'análises de líquidos biológicos (urina, sangue, líquor, etc.)', 'análises de tecidos', 'tratamentos alternativos ou complementares', 'acompanhamento nutricional', 'acompanhamento fisioterapêutico', 'acompanhamento fonoaudiológico', 'viagens (para avaliação de risco epidemiológico)', 'exposição a doenças contagiosas', 'doações de órgãos ou tecidos', 'vacinas recebidas (além das obrigatórias)', 'testes genéticos', 'injeções (vacinas, tratamentos, etc.)', 'cirurgias estéticas prévias', 'uso de dispositivos médicos (próteses, marcapassos, etc.)', 'problemas de visão', 'problemas de audição', 'problemas de mobilidade']
        self.political_preferences = ['afiliado político', 'afiliado politico', 'afiliados políticos', 'afiliados politicos', 'afiliação política', 'afiliaçao politica', 'alinhamento político', 'alinhamento politico', 'alinhamentos políticos', 'alinhamentos politicos', 'anarquismo', 'anarquista', 'anarquistas', 'associação política', 'associaçao politica', 'associações políticas', 'associaçoes politicas', 'bolsonarista', 'bolsonaristas', 'centrismo', 'centrista', 'centristas', 'comunismo', 'comunista', 'comunistas', 'conservador', 'conservadores', 'conservadorismo', 'de direita', 'de esquerda', 'direitista', 'direitismo', 'esquerdista', 'esquerdismo', 'extrema direita', 'extrema esquerda', 'filiação partidária', 'filiaçao partidaria', 'filiações partidárias', 'filiaçoes partidarias', 'ideologia política', 'ideologias políticas', 'ideologia politica', 'ideologias politicas', 'independente', 'independente', 'independentes', 'liberais', 'liberal', 'liberalismo', 'opiniao política', 'opiniões políticas', 'partido político', 'partidos políticos', 'Partido dos Trabalhadores', 'Partido da Social Democracia Brasileira', 'Partido do Movimento Democrático Brasileiro', 'Partido Democrático Trabalhista', 'Partido Socialista Brasileiro', 'Partido Comunista do Brasil', 'Partido Republicano Brasileiro', 'Partido Liberal', 'Partido Democratas', 'Partido Socialismo e Liberdade', 'Partido Verde', 'Partido Progressista', 'Partido Social Cristão', 'Partido Trabalhista Brasileiro', 'Partido da República', 'Partido Solidariedade', 'Partido Popular Socialista', 'Partido da Mobilização Nacional', 'Partido Social Democrata Cristão', 'Partido Trabalhista Cristão', 'Partido da Causa Operária', 'Partido Socialista dos Trabalhadores Unificado', 'Partido Humanista da Solidariedade', 'Partido Renovador Trabalhista Brasileiro', 'Partido Republicano Progressista', 'Partido da Mulher Brasileira', 'petista', 'petistas', 'posição política', 'posições políticas', 'progressismo', 'progressista', 'progressistas', 'radicais', 'radical', 'radicalismo', 'socialismo', 'socialista', 'socialistas', 'vinculação política', 'vinculações políticas', 'Afiliado a partido político', 'Filiação partidária', 'Ativista político', 'Engajamento político', 'Simpatizante político', 'Ideário político', 'Tendência política', 'Engajamento partidário', 'Política partidária', 'Ideologia partidária', 'Ideologia política', 'Voto político', 'Preferência política', 'Inclinação política', 'Opinião partidária', 'Inserção política', 'Engajamento cívico', 'Ativismo partidário', 'Inclinação ideológica', 'Atuação política', 'Aderência partidária', 'Adesão política', 'Linha política', 'Afinação política', 'Identificação partidária', 'Simpatia partidária', 'Filiação política', 'Visão política', 'Filiação ideológica', 'Posicionamento partidário', 'Posicionamento ideológico', 'Preferência partidária', 'Engajamento ideológico', 'Comprometimento político', 'Aderência ideológica', 'Tendência partidária', 'Identidade partidária', 'Identidade ideológica', 'Orientação política', 'Inclinação partidária']
        self.property_information = ['Histórico de propriedades', 'Endereço da propriedade', 'Tipo de propriedade', 'Valor da propriedade', 'Data de aquisição', 'Financiamentos imobiliários', 'Histórico de aluguel de imóveis', 'Endereço do imóvel alugado', 'Valor do aluguel', 'Metragem quadrada', 'Número de quartos', 'Número de banheiros', 'Ano de construção', 'Estado da propriedade', 'Avaliação imobiliária', 'Registro de propriedade', 'Descrição da propriedade', 'Histórico de vendas', 'Preço de venda', 'Data de venda', 'Histórico de hipotecas', 'Valor da hipoteca', 'Instituição financeira da hipoteca', 'Termos da hipoteca', 'Pagamento mensal da hipoteca', 'Datas de pagamento da hipoteca', 'Data de vencimento da hipoteca', 'Documentação imobiliária', 'Plano de zoneamento', 'Impostos sobre a propriedade', 'Histórico de impostos pagos', 'Regime de condomínio', 'Taxas de condomínio', 'Associação de proprietários', 'Restrições de propriedade', 'Easements (servidões)', 'Liens (ônus)', 'Notas fiscais', 'Contratos de locação', 'Termos de locação', 'Documentação do inquilino', 'Data de início do contrato de locação', 'Data de término do contrato de locação', 'Renovação de locação', 'Condições do imóvel alugado', 'Registros de reparos e manutenção', 'Seguro de propriedade', 'Cobertura do seguro', 'Detalhes de reclamações de seguro', 'Histórico de sinistros', 'Certidão de propriedade', 'Registro de escritura', 'Registro de matrícula', 'Registro de cadastro imobiliário', 'Laudo de vistoria', 'Relatório de inspeção de imóvel', 'Avaliação de mercado', 'Avaliação técnica', 'Avaliação de desempenho energético', 'Avaliação estrutural', 'Certificado de conformidade', 'Certificado de ocupação', 'Planta baixa', 'Planta arquitetônica', 'Planta hidráulica', 'Planta elétrica', 'Planta de gás', 'Planta de prevenção contra incêndios', 'Planta de acessibilidade', 'Planta de paisagismo', 'Regras do condomínio', 'Regras de uso da propriedade', 'Normas de construção', 'Restrições de zoneamento', 'Licenças de construção', 'Termos de licenciamento ambiental', 'Certificados de ocupação de solo', 'Limites de propriedade', 'Descrição de lote', 'Topografia do terreno', 'Declaração de propriedade rural', 'Dados de terras arrendadas', 'Informações de arrendamento agrícola', 'Registros de cultivo', 'Informações sobre pastagens', 'Registros de produção agrícola', 'Declarações de imposto de propriedade', 'Avaliações de propriedade rural', 'Dados de terras florestais', 'Licenças de corte de árvores', 'Registros de manejo florestal', 'Mapas de propriedade', 'Limites de propriedade', 'Descrição de lote', 'Marcas de propriedade', 'Cercas de propriedade', 'Acesso à propriedade', 'Easements (servidões)', 'Direitos de passagem', 'Conservação de terras', 'Restrições de conservação', 'Programas de conservação de terras', 'Acesso a recursos naturais', 'Registros de água', 'Direitos de água', 'Acesso a poços', 'Direitos de irrigação', 'Direitos de pesca', 'Direitos de caça', 'Direitos de mineração', 'Contratos de arrendamento mineral', 'Direitos de exploração de petróleo e gás', 'Direitos de energia renovável', 'Contratos de arrendamento de energia eólica', 'Contratos de arrendamento solar', 'Contratos de arrendamento hidrelétrico', 'Documentação de propriedade comercial', 'Contratos de locação comercial', 'Termos de locação comercial', 'Licenças comerciais']
        self.religion_keywords = ['cristianismo', 'cristão', 'cristã', 'cristãos', 'cristãs', 'cristao', 'cristaos', 'evangelismo', 'evangélico', 'evangélica', 'evangélicos', 'evangélicas', 'evangelico', 'evangelica', 'evangelicos', 'evangelicas', 'islã', 'islamismo', 'muçulmano', 'muçulmana', 'muçulmanos', 'muçulmanas', 'judaísmo', 'judaismo', 'judeu', 'judia', 'judeus', 'judias', 'judaico', 'judaica', 'judaicos', 'judaicas', 'hinduísmo', 'hinduismo', 'hindu', 'hindus', 'budismo', 'budista', 'budistas', 'siquismo', 'sikh', 'sikhs', 'taoísmo', 'taoísta', 'taoístas', 'taoismo', 'taoista', 'taoistas', 'xintoísmo', 'xintoísta', 'xintoístas', 'xintoismo', 'xintoista', 'xintoistas', 'espiritismo', 'espírita', 'espíritas', 'espirita', 'espiritas', 'candomblé', 'candomblecista', 'candomblecistas', 'umbanda', 'umbandista', 'umbandistas', 'espiritismo', 'paganismo', 'pagão', 'pagã', 'pagãos', 'pagãs', 'wicca', 'wiccano', 'wiccana', 'wiccanos', 'wiccanas', 'comunhão', 'culto', 'Santa Ceia', 'Mesquita', 'Ramadan', 'Sinagoga', 'Centro espírita', 'Desenvolvimento mediúnico', 'Terreiro', 'Ritual', 'Sabbat']
        self.travel_information = ['Viagens', 'Viagens internacionais', 'Viagens domésticas', 'Países visitados', 'Duração da viagem', 'Motivo da viagem', 'Viagens anteriores', 'Destinos visitados', 'Detalhes de hospedagem', 'Milhas acumuladas', 'Categoria de viajante', 'Tipo de bilhete', 'Classe da passagem', 'Assento', 'Aeroporto de partida', 'Aeroporto de chegada', 'Conexões de voo', 'Número do voo', 'Horário de partida', 'Horário de chegada', 'Compahia aérea', 'Código de reserva', 'Check-in online', 'Bagagem despachada', 'Bagagem de mão', 'Documentos de viagem', 'Passaporte', 'Visto de entrada', 'Seguro de viagem', 'Cartão de embarque', 'Imigração', 'Alfândega', 'Taxas de embarque', 'Taxas de segurança', 'Taxas de visto', 'Taxas de turismo', 'Vacinas obrigatórias', 'Aeroporto de trânsito', 'Aeroporto de conexão', 'Serviços adicionais', 'Aluguel de carro', 'Transfer do aeroporto', 'Wi-Fi no avião', 'Entretenimento a bordo', 'Refeições servidas', 'Política de cancelamento', 'Política de reembolso', 'Alterações de voo', 'Atrasos de voo', 'Cancelamentos de voo', 'Reacomodação de voo', 'Condições meteorológicas', 'Traslado para o hotel', 'Transporte público', 'Táxi/Uber/Lyft', 'Locais turísticos', 'Passeios guiados', 'Excursões', 'Cartão de turista', 'Mapas da cidade', 'Guia turístico', 'Segurança do destino', 'Clima do destino', 'Moeda local', 'Taxa de câmbio', 'Cartões de crédito aceitos', 'Idioma local', 'Cultura local', 'Gastronomia local', 'Compras', 'Artes e cultura', 'Vida noturna', 'Atrações turísticas', 'Museus e galerias', 'Parques e áreas naturais', 'Praias e balneários', 'Montanhas e trilhas', 'Eventos locais', 'Festivais e celebrações', 'Transporte urbano', 'Aluguel de bicicletas', 'Passeios de barco', 'Atividades ao ar livre', 'Esportes aquáticos', 'Aventuras radicais', 'Spa e bem-estar', 'Fitness e academia', 'Compras duty-free', 'Wi-Fi gratuito', 'Carregadores de celular', 'Áreas de descanso', 'Banheiros e chuveiros', 'Área de recreação infantil', 'Lojas e boutiques', 'Restaurantes e cafés', 'Balcões de informação', 'Salas VIP', 'Centro de negócios', 'Sala de conferências', 'Serviços médicos', 'Assistência especial']
        self.vehicle_information = ['registro de veículo', 'Placa do veículo', 'Marca do veículo', 'Modelo do veículo', 'Cor do veículo', 'Número do chassi', 'Número do motor', 'Ano do veículo', 'Infrações de trânsito', 'Tipo de infração', 'Local da infração', 'Data da infração', 'Carteira Nacional de Habilitação', 'Número do Renavam', 'Categoria da CNH', 'Validade da CNH', 'Pontuação na CNH', 'Multas de trânsito', 'Taxa de licenciamento', 'Seguro do veículo', 'Data de compra do veículo', 'Histórico de acidentes', 'Histórico de manutenção', 'Endereço de registro do veículo', 'Número do registro de motorista', 'Histórico de revisões', 'Histórico de recalls', 'Restrições judiciais no veículo', 'Número de série do veículo', 'Documento de transferência de veículo', 'Número de registro do veículo antigo', 'Data de vencimento do seguro', 'Teste de emissão de gases', 'Inspeção veicular', 'Certificado de registro de veículo', 'Autorização para conduzir ciclomotor', 'Autorização para conduzir veículo automotor', 'Autorização para conduzir veículo de transporte de carga', 'Autorização para conduzir veículo de transporte de passageiro', 'Autorização para conduzir veículo de transporte de escolares', 'Documento único do veículo', 'Certidão negativa de propriedade de veículo', 'Restrições administrativas no veículo', 'Restrições financeiras no veículo', 'Certificado de segurança veicular', 'Autorização especial de trânsito', 'Certificado de conclusão de curso de formação de condutores', 'Declaração de opção pelo regime de tributação especial', 'Comprovante de inscrição no cadastro de pessoas físicas', 'Comprovante de inscrição no cadastro nacional de pessoas jurídicas', 'Certificado de registro e licenciamento de veículo', 'Certificado de registro de veículo', 'Número de registro de identificação de veículo', 'Autorização para estacionamento em vaga reservada', 'Certificado de segurança veicular', 'Identificação do responsável pelo transporte', 'Certificado de identificação veicular', 'Número do protótipo', 'Documento fiscal eletrônico de prestação de serviço de transporte', 'Informações de registro de veículo em sistema nacional', 'Registro nacional de transportadores rodoviários de cargaDocumento fiscal eletrônico de transporte (DFT)', 'Relatório de inspeção de segurança veicular', 'Comprovante de seguro obrigatório']

        self.check_arguments()

    def print_help_and_exit(self):
        parser.print_help()
        sys.exit()

    def validate_directory_and_file_args(self):
        if not (self.args.directory or self.args.file):
            raise ValueError(f"[!] Erro: As opções -D ou -F devem ser fornecidas.")

        if self.args.directory and self.args.file:
            raise ValueError(f"[!] Erro: As opções -D e -F não podem ser fornecidas simultaneamente.")

    def validate_file_operations(self):
        if self.args.copy and (self.args.move or self.args.delete):
            raise ValueError(f"[!] Erro: --copy-files e (--move-files, --delete-files) não podem ser fornecidos simultaneamente.")

        if self.args.move and (self.args.copy or self.args.delete):
            raise ValueError(f"[!] Erro: --move-files e (--copy-files, --delete-files) não podem ser fornecidos simultaneamente.")

        if self.args.delete and (self.args.copy or self.args.move):
            raise ValueError(f"[!] Erro: --delete-files e (--copy-files, --move-files) não podem ser fornecidos simultaneamente.")

        if self.args.delete and (self.args.loot or self.args.database):
            raise ValueError(f"[!] Erro: --delete-files e (--loot, --database) não podem ser fornecidos simultaneamente.")

        if self.args.move and (self.args.loot or self.args.database):
            raise ValueError(f"[!] Erro: --move-files e (--loot, --database) não podem ser fornecidos simultaneamente.")

    def validate_findme_and_datatype_args(self):
        if self.args.findme and self.args.datatype:
            raise ValueError(f"[!] Erro: As opções --find e --data-type não podem ser fornecidas simultaneamente.")

    def validate_files(self):
        if self.args.file:
            self.args.file = self.validate_paths(self.args.file, is_dir=False)

    def validate_directories(self):
        if self.args.directory:
            self.args.directory = self.validate_paths(self.args.directory, is_dir=True)

    def validate_paths(self, paths, is_dir=False):
        validated_paths = set()
        for path in paths:
            if not os.access(path, os.R_OK):
                raise PermissionError(f"[!] Erro: Permissões insuficientes para ler {path}.")
            if not os.path.exists(path):
                raise FileNotFoundError(f'[!] Erro: O {"Diretório" if is_dir else "Arquivo"} fornecido {path} não existe.')
            if is_dir and os.path.isfile(path):
                raise ValueError(f"[!] Erro: O caminho fornecido {path} é um arquivo.")
            if not is_dir and os.path.isdir(path):
                raise ValueError(f"[!] Erro: O caminho fornecido {path} é um diretório.")
            validated_paths.add(path)
        return validated_paths

    def process_findme(self):
        if self.args.findme:
            self.args.findme = set(self.args.findme)
            self.regex_patterns.clear()
            for value in self.args.findme:
                value_key = ''.join(value.split())
                self.regex_patterns[f'{value_key}'] = re.compile(re.escape(value), re.IGNORECASE)

    def validate_tesseract(self):
        if platform.system().lower() in ['linux', 'linux2', 'darwin']:
            cmd = os.popen('which tesseract').read().strip()
            if not cmd:
                raise FileNotFoundError(f"[!] Não foi possível localizar o executável do Tesseract OCR.")
            pytesseract.tesseract_cmd = cmd
        elif platform.system().lower() in ['win32', 'cygwin', 'windows']:
            default_path = os.path.join(os.environ.get('ProgramFiles', 'C:\\Program Files'), 'Tesseract-OCR\\tesseract.exe')
            if os.path.exists(default_path):
                pytesseract.tesseract_cmd = default_path
            else:
                cmd = 'powershell -Command "(Get-Command -Name tesseract | Select-Object -ExpandProperty Path) -or (where.exe tesseract)"'
                path = os.popen(cmd).read().strip()
                if not path:
                    raise FileNotFoundError(f"[!] Não foi possível localizar o executável do Tesseract OCR.")
                pytesseract.tesseract_cmd = path

    def validate_datatype_args(self):
        if self.args.datatype:
            self.args.datatype = set(self.args.datatype)
            for dt in self.args.datatype:
                if dt not in self.regex_patterns:
                    raise ValueError(f"[!] Erro: Tipos de dados inválidos fornecidos: {dt}")
            for key in list(self.regex_patterns.keys()):
                if key not in self.args.datatype:
                    del self.regex_patterns[key]

    def validate_filetype_args(self):
        if self.args.filetype:
            self.args.filetype = set(self.args.filetype)
            for ft in self.args.filetype:
                if ft not in self.file_patterns:
                    raise ValueError(f"[!] Erro: Tipos de arquivo inválidos fornecidos: {ft}")
            for key in list(self.file_patterns.keys()):
                if key not in self.args.filetype:
                    del self.file_patterns[key]

    def validate_loot(self):
        if self.args.loot:
            if len(self.args.loot) > 1:
                raise ValueError(f"[!] Erro: O argumento --loot não pode ter vários valores.")
            else:
                self.args.loot = os.path.abspath(self.args.loot[0])
                if os.access(self.args.loot, os.W_OK):
                    if os.path.exists(f'{self.args.loot}/loot'):
                        raise ValueError(f"[!] Erro: {f'{self.args.loot}/loot'} já existe")
                    self.args.loot = f'{self.args.loot}/loot'
                else:
                    raise PermissionError(f"[!] Erro: Permissões insuficientes para escrever em {self.args.loot}")

    def validate_database(self):
        if self.args.database:
            if len(self.args.database) > 1:
                raise ValueError(f"[!] Erro: O argumento --database não pode ter vários valores.")
            else:
                self.args.database = os.path.abspath(self.args.database[0])
                if self.args.database != os.getcwd():
                    db_path = os.path.split(self.args.database)[0]
                    db_name = os.path.split(self.args.database)[1] if os.path.split(self.args.database)[1].split('.')[-1].lower() == 'db' else f'{os.path.split(self.args.database)[1]}.db'
                    self.args.database = os.path.join(db_path, db_name)
                else:
                    db_path = self.args.database
                    db_name = f'{platform.node()}.db'
                    self.args.database = os.path.join(db_path, db_name)
                if os.access(db_path, os.W_OK):
                    if os.path.exists(db_name):
                        raise FileExistsError(f"[!] Erro: {os.path.abspath(self.args.database)} já existe")
                else:
                    raise PermissionError(f"[!] Erro:Permissões insuficientes para escrever em {db_path}")

    def validate_copy(self):
        if self.args.copy:
            if os.path.exists(self.args.copy):
                if not os.access(self.args.copy, os.W_OK):
                    raise PermissionError(f"[!] Erro: Permissões insuficientes para escrever em {self.args.copy}")
            else:
                raise FileNotFoundError(f"[!] Erro: O caminho de destino {self.args.copy} não existe")

    def validate_move(self):
        if self.args.move:
            if os.path.exists(self.args.move):
                if not os.access(self.args.move, os.W_OK):
                    raise PermissionError(f"[!] Erro: Permissões insuficientes para escrever em {self.args.move}")
            else:
                raise FileNotFoundError(f"[!] Erro: O caminho de destino: {self.args.move} não existe")

    def check_arguments(self):
        if len(sys.argv) == 1:
            self.print_help_and_exit()

        try:
            self.validate_findme_and_datatype_args()
            self.validate_directory_and_file_args()
            self.validate_file_operations()
            self.validate_datatype_args()
            self.validate_filetype_args()
            self.validate_directories()
            self.validate_tesseract()
            self.validate_database()
            self.validate_files()
            self.validate_move()
            self.validate_loot()
            self.validate_copy()
            self.process_findme()

        except (ValueError, FileNotFoundError, FileExistsError, PermissionError) as e:
            print(e)
            self.print_help_and_exit()

        else:
            print(f"Iniciando SearchParty com PID {os.getpid()}. [{time.strftime('%H:%M:%S %d/%b/%Y')}]\n")
            self.start_information()

    def convert_size(self, size):
        for unit in ['Bytes', 'Kb', 'Mb', 'Gb', 'Tb']:
            if size < 1000.0:
                return f"{size:.2f} {unit}"
            size /= 1000.0

    def get_filesystem_size(self, path=os.path.abspath(os.sep)):
        total_size = psutil.disk_usage(path).total
        return self.convert_size(total_size)

    def start_information(self):

        current_user = os.getlogin()
        hostname = platform.node()
        fs_size = self.get_filesystem_size()
        operating_system = f'{platform.system()} {platform.release()} {platform.machine()}'
        sys_uptime = str(datetime.datetime.now() - datetime.datetime.fromtimestamp(psutil.boot_time())).split('.')[0]

        self.sys_info = f"[*] Informações do sistema\n\n" \
                        f"Usuário: {current_user}\n" \
                        f"Máquina: {hostname}\n" \
                        f"Sistema operacional: {operating_system}\n" \
                        f"Sistema de arquivos: {fs_size}\n" \
                        f"Tempo de atividade: {sys_uptime}\n"

        if self.args.directory:
            self.sys_info += f"Diretórios: {' '.join(directory for directory in self.args.directory)}"

        if self.args.file:
            self.sys_info += f"Aqruivos: {' '.join(file_path for file_path in self.args.file)}"

        print(self.sys_info)
        self.process_contents()

    def process_contents(self):

        self.supported_files = {}
        print(f"\n[*] Enumerando conteúdos. [{time.strftime('%H:%M:%S %d/%b/%Y')}]\n")
        if self.args.directory:
            self.directory_size = 0
            for argument in self.args.directory:
                for path, dirs, files in os.walk(argument):
                    for filename in files:
                        full_path = os.path.abspath(os.path.join(path, filename))
                        self.directory_size += os.path.getsize(full_path)
                        for pattern_name, mime_type in self.file_patterns.items():
                            if from_file(full_path, mime=True) == mime_type or guess_mime(full_path) == mime_type:
                                if pattern_name not in self.supported_files:
                                    self.supported_files[pattern_name] = []
                                self.supported_files[pattern_name].append(full_path)

        if len(self.supported_files) != 0:
            print(f"Os diretórios possuem {sum(len(value) for value in self.supported_files.values())} arquivos suportados. {self.convert_size(self.directory_size)}")
            sorted_keys = sorted(self.supported_files.keys(), key=lambda x: (self.file_patterns[x], x, len(x)))
            for key in sorted_keys:
                print(f"{self.file_patterns[key]} ({len(self.supported_files[key])})")

        if self.args.file:
            for argument in self.args.file:
                for pattern_name, mime_type in self.file_patterns.items():
                    if from_file(argument, mime=True) == mime_type or guess_mime(argument) == mime_type:
                        if pattern_name not in self.supported_files:
                            self.supported_files[pattern_name] = []
                        self.supported_files[pattern_name].append(argument)
            if len(self.supported_files) != 0:
                for value in self.supported_files.values():
                    for file_path in value:
                        print(f"File: {file_path} ({self.convert_size(os.path.getsize(file_path))}) is {guess_mime(file_path) or from_file(file_path, mime=True)}")

        if len(self.supported_files) == 0:
            print(f"[-] Nenhum arquivos suportado foi encontrado")
            sys.exit()

        self.process_files()

    def process_files(self):

        print(f"\n[*] Iniciando a extração de dados. [{time.strftime('%H:%M:%S %d/%b/%Y')}]\n")

        self.data_found = {}
        self.error_files = {}
        self.data_indexes = {}

        extraction_methods = {
            'bmp': self.extract_from_bmp, 'csv': self.extract_from_csv, 'gif': self.extract_from_gif,
            'pdf': self.extract_from_pdf, 'png': self.extract_from_png, 'txt': self.extract_from_txt,
            'docx': self.extract_from_docx, 'jpeg': self.extract_from_jpeg, 'pptx': self.extract_from_pptx,
            'tiff': self.extract_from_tiff, 'webp': self.extract_from_webp, 'xlsx': self.extract_from_xlsx,
            'email': self.extract_from_email
        }

        for file_type, method in extraction_methods.items():
            if self.supported_files.get(file_type) is not None:
                method(self.supported_files[file_type])

        self.process_results()

    def extract_from_pdf(self, pdf_path):
        print("Extraindo texto de arquivos PDF")
        text = ''
        try:
            for pdf_file in pdf_path:
                text = extract_text(pdf_file)
                self.find_patterns(text, pdf_file)
        except Exception:
            self.error_files['pdf'] = self.error_files.get('pdf', 0) + 1

    def extract_from_pptx(self, pptx_path):
        print("Extraindo texto de arquivos PowerPoint")
        for pptx_file in pptx_path:
            try:
                text = ''
                presentation = Presentation(pptx_file)
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            for line in shape.text.split('\n'):
                                text += f'{line}\n'
                self.find_patterns(text, pptx_file)
            except Exception:
                if 'pptx' not in self.error_files:
                    self.error_files['pptx'] = 1
                self.error_files['pptx'] += 1

    def extract_from_docx(self, docx_path):
        print("Extraindo texto de documentos Word")
        for docx_file in docx_path:
            try:
                text = ''
                wdoc = Document(docx_file)
                for wdoc_p in wdoc.paragraphs:
                    for line in wdoc_p.text.split('\n'):
                        text += f'{line}\n'
                self.find_patterns(text, docx_file)
            except Exception:
                if 'docx' not in self.error_files:
                    self.error_files['docx'] = 1
                self.error_files['docx'] += 1

    def extract_from_xlsx(self, xlsx_path):
        print("Extraindo texto de SpreadSheets")
        for xlsx_file in xlsx_path:
            try:
                text = ''
                wbook = load_workbook(filename=xlsx_file)
                wsheet = wbook.active
                for row in wsheet.iter_rows(values_only=True):
                    for cell in row:
                        for line in str(cell).split('\n'):
                            text += f'{line}\n'
                self.find_patterns(text, xlsx_file)
            except Exception:
                if 'xlsx' not in self.error_files:
                    self.error_files['xlsx'] = 1
                self.error_files['xlsx'] += 1

    def extract_from_txt(self, txt_path):
        print("Extraindo texto de arquivos de Texto")
        for txt_file in txt_path:
            try:
                with open(txt_file, 'r', encoding='utf-8') as file:
                    text = file.read()
                    self.find_patterns(text.strip(), txt_file)
            except Exception:
                if 'txt' not in self.error_files:
                    self.error_files['txt'] = 1
                self.error_files['txt'] += 1

    def extract_from_csv(self, csv_path):
        print("Extraindo texto de arquivos CSV")
        for csv_file in csv_path:
            try:
                text = ''
                with open(csv_file, 'r', encoding='utf-8') as file:
                    reader = csv.reader(file)
                    for row in reader:
                        text += str(row)
                self.find_patterns(text, csv_file)
            except Exception:
                if 'csv' not in self.error_files:
                    self.error_files['csv'] = 1
                self.error_files['csv'] += 1

    def extract_from_jpeg(self, jpeg_path):
        print("Extraindo texto de imagens JPEG")
        for jpeg_file in jpeg_path:
            try:
                text = ''
                for line in pytesseract.image_to_string(jpeg_file).split('\n'):
                    text += f'{line}\n'
                self.find_patterns(text, jpeg_file)
            except Exception:
                if 'jpeg' not in self.error_files:
                    self.error_files['jpeg'] = 1
                self.error_files['jpeg'] += 1

    def extract_from_png(self, png_path):
        print("Extraindo texto de imagens PNG")
        for png_file in png_path:
            try:
                text = ''
                for line in pytesseract.image_to_string(png_file).split('\n'):
                    text += f'{line}\n'
                self.find_patterns(text, png_file)
            except Exception:
                if 'png' not in self.error_files:
                    self.error_files['png'] = 1
                self.error_files['png'] += 1

    def extract_from_bmp(self, bmp_path):
        print("Extraindo texto de imagens Bitmap")
        for bmp_file in bmp_path:
            try:
                text = ''
                for line in pytesseract.image_to_string(bmp_file).split('\n'):
                    text += f'{line}\n'
                self.find_patterns(text, bmp_file)
            except Exception:
                if 'bmp' not in self.error_files:
                    self.error_files['bmp'] = 1
                self.error_files['bmp'] += 1

    def extract_from_tiff(self, tiff_path):
        print("Extraindo texto de imagens TIFF")
        for tiff_file in tiff_path:
            try:
                text = ''
                for line in pytesseract.image_to_string(tiff_file).split('\n'):
                    text += f'{line}\n'
                self.find_patterns(text, tiff_file)
            except Exception:
                if 'tiff' not in self.error_files:
                    self.error_files['tiff'] = 1
                self.error_files['tiff'] += 1

    def extract_from_webp(self, webp_path):
        print("Extraindo texto de imagens WebP")
        for webp_file in webp_path:
            try:
                text = ''
                for line in pytesseract.image_to_string(webp_file).split('\n'):
                    text += f'{line}\n'
                self.find_patterns(text, webp_file)
            except Exception:
                if 'webp' not in self.error_files:
                    self.error_files['webp'] = 1
                self.error_files['webp'] += 1

    def extract_from_gif(self, gif_path):
        print("Extraindo texto de imagens GIF")
        for gif_file in gif_path:
            try:
                with Image.open(gif_file) as file:
                    extracted_set = set()
                    for frame in range(file.n_frames):
                        file.seek(frame)
                        frame_image = file.copy()
                        frame_text = pytesseract.image_to_string(frame_image)
                        for text in frame_text.split('\n'):
                            if text not in extracted_set:
                                extracted_set.add(text)
                    self.find_patterns(text, gif_file)
            except Exception:
                if 'gif' not in self.error_files:
                    self.error_files['gif'] = 1
                self.error_files['gif'] += 1

    def extract_from_email(self, email_path):
        print(f"Extraindo texto de arquivos de Email")
        for email_file in email_path:
            try:
                with open(email_file, 'r', encoding='utf-8') as file:
                    msg = message_from_file(file)
                    text = ''
                    for header in [msg[header] for header in msg.keys() if header.lower() in ['from', 'to', 'subject']]:
                        text += f'{header}\n'
                    if msg.is_multipart():
                        for part in msg.walk():
                            if part.get_content_type() == 'text/plain':
                                payload = part.get_payload(decode=True)
                                for line in payload.splitlines():
                                    line = line.decode(part.get_content_charset()).split("\n")
                                    text += f'{line}\n'
                        self.find_patterns(text, email_file)
                    else:
                        payload = msg.get_payload(decode=True)
                        for line in payload.splitlines():
                            text += f'{line.decode()}\n'
                        self.find_patterns(text, email_file)
            except Exception:
                if 'emails' not in self.error_files:
                    self.error_files['emails'] = 1
                self.error_files['emails'] += 1

    def find_patterns(self, text, path):
        keywords_found = {'Grupos étnicos': [keyword for keyword in self.ethnic_groups if re.search(keyword, text, re.IGNORECASE)],
                          'Informações financeiras': [keyword for keyword in self.financial_information if re.search(keyword, text, re.IGNORECASE)],
                          'Gênero e orientação sexual': [keyword for keyword in self.gender_orientation if re.search(keyword, text, re.IGNORECASE)],
                          'Informações jurídicas': [keyword for keyword in self.legal_information if re.search(keyword, text, re.IGNORECASE)],
                          'Informações médicas': [keyword for keyword in self.medical_information if re.search(keyword, text, re.IGNORECASE)],
                          'Preferências políticas': [keyword for keyword in self.political_preferences if re.search(keyword, text, re.IGNORECASE)],
                          'Informações sobre imóveis': [keyword for keyword in self.property_information if re.search(keyword, text, re.IGNORECASE)],
                          'Religião e fé': [keyword for keyword in self.religion_keywords if re.search(keyword, text, re.IGNORECASE)],
                          'Histórico de viagens': [keyword for keyword in self.travel_information if re.search(keyword, text, re.IGNORECASE)],
                          'Documentos de veículo': [keyword for keyword in self.vehicle_information if re.search(keyword, text, re.IGNORECASE)]}

        if self.args.findme:
            for argument in self.args.findme:
                key = ''.join(argument.split())
                if any(value.search(text) for value in [self.regex_patterns[key]]):
                    if argument not in self.data_found:
                        self.data_found[argument] = []
                    info = {'file': path}
                    if any(value for value in keywords_found.values()):
                        info['keywords'] = keywords_found
                    self.data_found[argument].append(info)
                    self.data_indexes[argument] = self.data_indexes.get(argument, 0) + 1
        else:
            for pattern_name, pattern in self.regex_patterns.items():
                for match in pattern.finditer(text):
                    self.data_indexes[pattern_name] = self.data_indexes.get(pattern_name, 0) + 1
                    if f"{pattern_name}" not in self.data_found:
                        self.data_found[f"{pattern_name}"] = set()
                    for key, value in ((key, value) for key, value in keywords_found.items() if value):
                        if key not in self.data_found:
                            self.data_found[key] = set()
                        self.data_indexes[key] = self.data_indexes.get(key, 0) + 1
                        self.data_found[key].add(path)
                    self.data_found[f"{pattern_name}"].add(path)

    def copy_files(self):
        print(f"[*] Copiando arquivos para {self.args.copy} por favor aguarde.")

        errors = set()
        files_copied = False

        for file_list in self.data_found.values():
            if file_list:
                for info in file_list:
                    try:
                        if self.args.findme:
                            if not os.path.exists(os.path.join(self.args.copy, os.path.basename(info['file']))):
                                shutil.copy(info['file'], self.args.copy)
                        else:
                            if not os.path.exists(os.path.join(self.args.move, os.path.basename(info))):
                                shutil.copy(info, self.args.copy)
                        files_copied = True
                    except Exception as error:
                        error_type = type(error).__name__
                        errors.add(error_type)

        copied_message = f'\n[+] Concluído!' if files_copied else ''

        if errors:
            copied_message += f'\n[-] Errors foram encontrados ao copiar arquivos: {", ".join(errors)}'

        print(copied_message)

    def move_files(self):
        print(f"[*] Movendo arquivos para {self.args.move} por favor aguarde.")

        errors = set()
        files_moved = False

        for data_type, file_list in self.data_found.items():
            if file_list:
                for info in file_list:
                    try:
                        if self.args.findme:
                            if not os.path.exists(os.path.join(self.args.move, os.path.basename(info['file']))):
                                shutil.move(info['file'], self.args.move)
                        else:
                            if not os.path.exists(os.path.join(self.args.move, os.path.basename(info))):
                                shutil.move(info, self.args.move)
                        files_moved = True
                    except Exception as error:
                        error_type = type(error).__name__
                        errors.add(error_type)

        moved_message = f'\n[+] Concluído!' if files_moved else ''

        if errors:
            moved_message += f'\nErrors foram encontrados ao mover arquivos: {", ".join(errors)}'

        print(moved_message)

    def delete_files(self):
        print(f"[*] Deletando arquivos, por favor aguarde.")

        errors = set()

        for data_type, file_list in self.data_found.items():
            if file_list:
                for info in file_list:
                    try:
                        if self.args.findme:
                            if os.path.exists(info['file']):
                                os.remove(info['file'])
                        else:
                            if os.path.exists(info):
                                os.remove(info)
                        files_deleted = True
                    except Exception as error:
                        error_type = type(error).__name__
                        errors.add(error_type)

        deleted_message = f'\n[+] Concluído!' if files_deleted else ''

        if errors:
            deleted_message += f'\n[-] Errors foram encontrados ao deletar arquivos: {", ".join(errors)}'

        print(deleted_message)

    def export_to_database(self):
        conn = sqlite3.connect(self.args.database)
        cursor = conn.cursor()
        for table_name, file_list in self.data_found.items():
            if file_list:
                cursor.execute(f'''CREATE TABLE "{table_name}" (files TEXT)''')
                for info in file_list:
                    if self.args.findme:
                        cursor.execute(f'''INSERT INTO "{table_name}" (files) VALUES (?)''', (info['file'],))
                    else:
                        cursor.execute(f'''INSERT INTO "{table_name}" (files) VALUES (?)''', (info,))
        conn.commit()
        conn.close()

    def export_to_loot(self):
        os.makedirs(self.args.loot)
        for data_type, file_list in self.data_found.items():
            if file_list:
                with open(f'{self.args.loot}/{data_type}.txt', 'w') as output_file:
                    if self.args.findme:
                        for info in file_list:
                            print(info['file'], file=output_file)
                    else:
                        for info in file_list:
                            print(info, file=output_file)

    def result_stats(self):
        results_info = f"\n[+] Análise de Dados:\n"
        results_info += ''.join([f"\n{key + ':': <26} {value} ocorrências" for key, value in self.data_indexes.items() if value != 0]) + '\n'
        results_info += f"\n[+] Análise de Arquivos:\n"

        if not (self.args.findme or self.args.file):
            results_info += ''.join([f"\n{key + ':': <27}{len(value)} arquivos" for key, value in self.data_found.items() if value != 0]) + '\n'

        results_info += ''.join([f"\n{str(key).upper() + ' arquivos:': <13}{len(self.supported_files[key])} processados" for key in self.supported_files.keys()])

        if sum(self.error_files.values()) != 0:
            results_info += f'\n\n[-] Erros foram encontrados ao ler {sum(self.error_files.values())} arquivos:\n'
            results_info += ''.join([f"\n{str(key).upper()+':': <6}{self.error_files[key]}" for key in self.error_files.keys()])

        results_info += f"\n\n[+] Processo de busca finalizado. [{time.strftime('%H:%M:%S %d/%b/%Y')}]"
        return results_info

    def results_information(self):

        if self.args.findme and (self.args.directory or self.args.file):
            arguments_not_found = []
            arguments_found = [argument for argument in self.args.findme if len(self.data_found[argument]) != 0]
            for argument in self.args.findme:
                if argument in arguments_found:
                    files_with_keywords = sum(1 for info in self.data_found[argument] if "keywords" in info)
                    print(f'\n[+] "{argument}" foi encontrado em {len(self.data_found[argument])} arquivo(s). {f"{files_with_keywords} contém PII/dados sensiveis." if files_with_keywords else ""}\n')
                    for info in self.data_found[argument]:
                        print(info['file'])
                    if files_with_keywords:
                        for info in [info for info in self.data_found[argument] if 'keywords' in info]:
                            print(f"\n[+] Palavras-chave encontradas em {info['file']}\n")
                            for key in (key for key in info['keywords'] if info['keywords'].get(key) and len(info['keywords'].get(key)) != 0):
                                print(f"{key}")
                else:
                    arguments_not_found.append(argument)
            if arguments_not_found:
                print()
                for argument in arguments_not_found:
                    print(f"[-] Não foi encontrada nenhuma correspondência para '{argument}' on {', '.join([os.path.basename(os.path.normpath(directory)) for directory in self.args.directory]) or ', '.join([os.path.basename(file_path) for file_path in self.args.file])}")
        else:
            if self.args.directory:
                for data_type, file_list in self.data_found.items():
                    if file_list:
                        print(f"\n[+] Arquivos contendo {data_type}:\n")
                        for info in file_list:
                            print(info)

    def process_results(self):
        if sum(value for value in self.data_indexes.values()) == 0:
            if self.args.directory:
                print(f"\n[-] Nenhum dado foi encontrado em: {', '.join([os.path.basename(os.path.abspath(directory)) for directory in self.args.directory])}\n")
            if self.args.file:
                print(f"\n[-] Nenhum dado foi encontrado em: {', '.join([os.path.basename(file_path) for file_path in self.args.file])}\n")
        else:

            if self.args.loot or self.args.database:
                print(f'{self.result_stats()}\n')
                (self.export_to_loot() if self.args.loot else None) or (self.export_to_database() if self.args.database else None)
                if self.args.loot and self.args.database:
                    print(f"Salvando resultados em {', '.join([self.args.loot, self.args.database])}")
                else:
                    print(f"Salvando resultados em {self.args.loot or self.args.database}" if self.args.loot or self.args.database else "")

            else:
                print(f"\n[!] Nenhum arquivo de saída especificado, mostrando os resultados no terminal.")
                time.sleep(5)
                self.results_information()
                print(f'{self.result_stats()}\n')
                if self.args.copy:
                    self.copy_files()
                if self.args.move:
                    self.move_files()
                if self.args.delete:
                    self.delete_files()


start_time = time.time()

if __name__ == '__main__':
    SearchParty()

end_time = time.time()

print(f'Tempo total de execução: {datetime.timedelta(seconds=end_time - start_time)}')
